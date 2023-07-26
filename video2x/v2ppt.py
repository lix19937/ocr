import pytube
from pathlib import Path
import argparse
import cv2
import os  
import numpy as np

import glob
from pptx import Presentation  
from pptx.util import Inches 
import pptx 
from pathlib import Path
import re

def natural_sort_key(s, _nsre=re.compile('([0-9]+)')):
    return [
        int(text)
        if text.isdigit() else text.lower()
        for text in _nsre.split(s)]

# import argparse
# ap = argparse.ArgumentParser()
# ap.add_argument("-p", "--path", required=False, default='', help="Path to load img files")
# args = vars(ap.parse_args())

def create_ppt(img_path, pptname):
    prs = Presentation()  
    blank_slide_layout = prs.slide_layouts[6]  

    img_path = Path(img_path)

    images = [img.name for img in img_path.iterdir() if img.suffix == ".jpg"]

    sorted_images = sorted(images, key=natural_sort_key)

    for img in sorted_images:
        slide = prs.slides.add_slide(blank_slide_layout) 
        pic = slide.shapes.add_picture(f"{img_path/Path(img)}", Inches(0.5), Inches(0.75), width=Inches(9.2), height=Inches(6)) 

    [f.unlink() for f in img_path.glob("*") if f.is_file()] 
    prs.save(img_path/Path(pptname + '.pptx'))



ap = argparse.ArgumentParser()
ap.add_argument("-s", "--start", required=False, type=int, default=0,	help="Start of the video")
ap.add_argument("-o", "--output", required=False, default='./output', help="Path to save slides")
ap.add_argument("-i", "--input", required=False, default='./input', help="Path of videos")

args = vars(ap.parse_args())

inpath = str(Path(args['input']))
outpath = str(Path(args['output']))

videos = glob.glob(inpath + "/*.mp4")   

for i in range(len(videos)):
  curr_video = videos[i].strip()
  pptname = os.path.basename(curr_video).split(".")[0] 
  print("pptname is|", pptname)

  vidcap = cv2.VideoCapture(f"{curr_video}")
  success, image = vidcap.read()
  print("img shape is ", image.shape)

  #################### Setting up parameters ################
  seconds = 1
  fps = vidcap.get(cv2.CAP_PROP_FPS) # Gets the frames per second
  multiplier = fps * seconds + 15    # 12 is just eliminate ghosting

  print("multiplier is ", multiplier)

  #################### Initiate Process ################
  count = 0
  threshold = 0.003
  real_save_cnt = 1
  stop_cnt = 10

  ####################  Run Process    ################
  while success:
      frameId = int(round(vidcap.get(1))) #current frame number, rounded b/c sometimes you get frame intervals which aren't integers...this adds a little imprecision but is likely good enough
      success, image = vidcap.read()
      if count == 0:
          previous = image
          # cv2.imwrite(f"{vid_path}/frame{frameId}.jpg", image)
      count =+ 1
      
      if real_save_cnt >= stop_cnt:
          break

      if frameId % int(multiplier) == 0 and frameId > int(multiplier)*args['start']:
          current = image
          try:
                      
              prev = cv2.cvtColor(previous, cv2.COLOR_BGR2GRAY)
              cur = cv2.cvtColor(current, cv2.COLOR_BGR2GRAY)

              (thresh, im_bw_prev) = cv2.threshold(prev, 128, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)
              (thresh, im_bw_cur) = cv2.threshold(cur, 128, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)

              # Set up and feed background subtractor (cf. tutorial linked in question)
              backSub = cv2.createBackgroundSubtractorMOG2()
              _ = backSub.apply(im_bw_prev)
              mask = backSub.apply(im_bw_cur)
              n_white_pix = np.sum(mask == 255)/mask.size

              if n_white_pix > threshold:
                  cv2.imwrite(f"{outpath}/frame{frameId}.jpg", image)
                  print("save ...{}".format(real_save_cnt))
                  real_save_cnt +=1
              
              previous = current
          except:
              continue

  vidcap.release()
  create_ppt(outpath, pptname)
  print("Complete")



## (v2ppt) lix19937@lix19937:~/Desktop/video2x$
# python v2ppt.py  -p ./input/
